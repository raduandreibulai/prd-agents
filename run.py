from typing import List
import csv
import os
import json
import sys
from pathlib import Path

from dotenv import load_dotenv
from openai import OpenAI

load_dotenv()
client = OpenAI()

MODEL = "gpt-4.1-mini"


def read_text(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


def read_prompt(path: Path) -> str:
    return path.read_text(encoding="utf-8")


def call_agent(prompt: str, user_input: str) -> str:
    resp = client.responses.create(
        model=MODEL,
        input=[
            {"role": "system", "content": prompt},
            {"role": "user", "content": user_input},
        ],
    )
    return resp.output_text

def ingest_txt_md(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")

def ingest_csv(path: Path) -> str:
    rows = []
    with path.open(newline="", encoding="utf-8", errors="ignore") as f:
        reader = csv.reader(f)
        for i, row in enumerate(reader):
            if i >= 30:  # limit preview
                break
            rows.append(row)
    return f"CSV PREVIEW (first 30 rows):\n{rows}"

def ingest_xlsx(path: Path) -> str:
    # lightweight summary (no heavy analysis)
    try:
        import openpyxl
    except ImportError:
        return "XLSX provided but openpyxl is not installed."

    wb = openpyxl.load_workbook(path, data_only=True)
    parts = []
    for sheet in wb.worksheets[:3]:  # limit to first 3 sheets
        parts.append(f"Sheet: {sheet.title}")
        # take a small top-left preview
        preview = []
        for r in range(1, 16):
            row_vals = []
            for c in range(1, 11):
                v = sheet.cell(row=r, column=c).value
                row_vals.append(v)
            preview.append(row_vals)
        parts.append(f"Top-left preview (15x10):\n{preview}")
    return "\n".join(parts)

def ingest_pptx_text_only(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        return "PPTX provided but python-pptx is not installed."

    prs = Presentation(str(path))
    parts = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                t = shape.text.strip()
                if t:
                    texts.append(t)
        notes_text = ""
        try:
            notes = slide.notes_slide
            if notes and notes.notes_text_frame:
                notes_text = notes.notes_text_frame.text.strip()
        except Exception:
            pass

        parts.append(f"--- Slide {i} ---")
        if texts:
            parts.append("SLIDE TEXT:\n" + "\n".join(texts))
        else:
            parts.append("SLIDE TEXT: (none detected)")

        if notes_text:
            parts.append("SPEAKER NOTES:\n" + notes_text)
        else:
            parts.append("SPEAKER NOTES: (none)")

        # Callout if slide likely contains images but no text
        if not texts and not notes_text:
            parts.append("NOTE: This slide may be mostly images/screenshots. Add speaker notes or a captions doc.")

    return "\n".join(parts)

def build_context_bundle(context_dir: Path) -> str:
    if not context_dir.exists():
        return ""

    chunks: List[str] = []
    for path in sorted(context_dir.glob("*")):
        if path.is_dir():
            continue
        ext = path.suffix.lower()
        header = f"\n\n===== FILE: {path.name} =====\n"
        if ext in [".txt", ".md"]:
            body = ingest_txt_md(path)
        elif ext == ".csv":
            body = ingest_csv(path)
        elif ext == ".xlsx":
            body = ingest_xlsx(path)
        elif ext == ".pptx":
            body = ingest_pptx_text_only(path)
        else:
            body = f"(Skipped unsupported file type: {ext})"
        chunks.append(header + body)

    return "\n".join(chunks).strip()

def main() -> None:
    if len(sys.argv) < 2:
        print("Usage: python run.py transcripts/meeting1.txt")
        raise SystemExit(1)

    transcript_path = Path(sys.argv[1])
    if not transcript_path.exists():
        print(f"ERROR: Transcript file not found: {transcript_path}")
        raise SystemExit(1)

    Path("out").mkdir(exist_ok=True)

    # Load prompts
    chief_prompt = read_prompt(Path("agents/chief_of_staff.txt"))
    prd_prompt = read_prompt(Path("agents/prd_writer.txt"))
    cfo_prompt = read_prompt(Path("agents/cfo_critic.txt"))
    coo_prompt = read_prompt(Path("agents/coo_jira.txt"))

    # Read transcript
    transcript = read_text(transcript_path)
    context_bundle = build_context_bundle(Path("context"))
    print("DEBUG: Context bundle length =", len(context_bundle))
    print("HIT: building context pack...")

    context_pack_md = ""
    if context_bundle.strip():
    	librarian_prompt = read_prompt(Path("agents/librarian.txt"))
    	context_pack_md = call_agent(librarian_prompt, context_bundle)
    	(Path("out") / "context.md").write_text(context_pack_md, encoding="utf-8")
    	print("HIT: wrote out/context.md")
    else:
    	print("HIT: context folder empty or unreadable")

    # Combine transcript + context pack for downstream agents
    combined_input = transcript
    if context_pack_md.strip():
    	combined_input = f"{transcript}\n\n===== CONTEXT PACK =====\n{context_pack_md}"
    	print("HIT: passing transcript + context pack to Chief of Staff")
    else:
    	print("HIT: passing transcript ONLY to Chief of Staff (no context pack)")

    librarian_prompt = read_prompt(Path("agents/librarian.txt")) if context_bundle else ""
    context_pack_md = ""
    if context_bundle:
        context_pack_md = call_agent(librarian_prompt, context_bundle)
        (Path("out") / "context.md").write_text(context_pack_md, encoding="utf-8")

    # Combine transcript + context pack for downstream agents
    combined_input = transcript
    if context_pack_md:
        combined_input = f"{transcript}\n\n===== CONTEXT PACK =====\n{context_pack_md}"

    # 1) Chief of Staff -> notes JSON (or raw)
    notes_text = call_agent(chief_prompt, combined_input)
    (Path("out") / "notes_raw.txt").write_text(notes_text, encoding="utf-8")

    notes_json = None
    try:
        notes_json = json.loads(notes_text)
        (Path("out") / "notes.json").write_text(json.dumps(notes_json, indent=2), encoding="utf-8")
    except json.JSONDecodeError:
        pass

    # 2) PRD Writer
    prd_input = json.dumps(notes_json, indent=2) if notes_json else notes_text
    prd_md = call_agent(prd_prompt, prd_input)
    (Path("out") / "prd.md").write_text(prd_md, encoding="utf-8")

    # 3) CFO Critic
    roi_review = call_agent(cfo_prompt, prd_md)
    (Path("out") / "roi_review.md").write_text(roi_review, encoding="utf-8")

    # 4) COO -> Jira JSON
    jira_json_text = call_agent(coo_prompt, prd_md)
    (Path("out") / "jira_plan.json").write_text(jira_json_text, encoding="utf-8")

    print("✅ Done!")
    print("Generated:")
    print("- out/notes_raw.txt")
    if notes_json:
        print("- out/notes.json")
    print("- out/prd.md")
    print("- out/roi_review.md")
    print("- out/jira_plan.json")


if __name__ == "__main__":
    main()
